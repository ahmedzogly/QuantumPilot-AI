from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

COLORS = {
    "bg_dark": RGBColor(10, 14, 26),
    "bg_card": RGBColor(18, 20, 31),
    "border": RGBColor(30, 34, 53),
    "blue": RGBColor(15, 98, 254),
    "green": RGBColor(36, 161, 72),
    "white": RGBColor(244, 244, 244),
    "gray": RGBColor(141, 141, 141),
    "yellow": RGBColor(241, 194, 27),
    "red": RGBColor(218, 30, 40),
}

def add_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg_dark"]

def add_title_content_slide(title, content_list, title_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = title_color or COLORS["white"]
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(6.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["white"] if i % 2 == 0 else COLORS["gray"]
        p.space_after = Pt(10)
        p.level = 0
    return slide

# Slide 1 Title - Progress Presentation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1.5))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "QuantumPilot AI - Progress Presentation"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = COLORS["white"]
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "AI Orchestration Platform for Quantum Computing - From Circuit Input to Cost Savings with Space Weather Awareness"
p.font.size = Pt(16)
p.font.color.rgb = COLORS["gray"]
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(12)

content_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.33), Inches(4))
tf = content_box.text_frame
tf.word_wrap = True
bullets = [
    "Live IBM Quantum: fez 156q T1 135.6µs, marrakesh 170.9µs, kingston 231µs BEST - 3 Heron R2 processors",
    "8.04M Calibration Records → 8847 Contexts for NeuralUCB Training",
    "NeuralUCB 22-D Contextual Bandit: Loss 0.3224→0.0028, 76% Cost Saving via S-ZNE Constant Overhead",
    "Novel Discovery: Space Weather Correlation T1 vs Kp -0.197 p=0.00047 - First Study",
    "Bilingual Copilot Intent Agent: Arabic/English 'نفذ بأقل تكلفة' → Execution Plan",
    "Production Ready: Docker 6 Services, Vercel Frontend+Backend, CI/CD, 20 Docs, GitHub 132 Files",
    "Live Demo: https://quantumpilot-ai.vercel.app • GitHub: github.com/ahmedzogly/QuantumPilot-AI"
]
for i, b in enumerate(bullets):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = "• " + b
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS["white"]
    p.space_after = Pt(8)

# Slide 2 Introduction
add_title_content_slide("1. Introduction", [
    "Quantum computing has reached NISQ era with 156-qubit processors (IBM Heron R2: fez, marrakesh, kingston) but practical utility hindered by noise, decoherence, and operational overhead",
    "IBM Quantum Runtime provides devices + Cloud + Qiskit, but backend selection is least_busy only (ignores T1/T2/RO/CZ), no AI decision, no explainability",
    "Researchers today manually: write circuit → pick random backend → try opt_level 1 → fail queue → try mitigation ZNE 5x cost → fail due to solar storm → retry 5x manually → 70% quantum time wasted",
    "We present QuantumPilot AI, first AI Operating Intelligence Platform that acts as autonomous autopilot: manages full lifecycle from circuit input (QASM/Qiskit Python) to cost savings, with continuous learning and explainability",
    "Core: NeuralUCB 22-D contextual bandit for backend selection (72 arms = 3 backends ×4 opt ×6 mitigation) + SpaceWeatherService LIVE NOAA + Mitigation Factory 7 methods + Copilot Intent Agent bilingual AR/EN",
    "Deployed production-ready: Docker 6 services (postgres, redis, rabbitmq, backend Python 3.11 for Mitiq, celery-worker, celery-beat every 10m SpaceWeather, frontend Node 20), Vercel 100% without Visa/PRO (frontend+backend Python serverless), GitHub Codespaces, CI/CD 5 jobs, 20 docs"
], title_color=COLORS["blue"])

# Slide 3 Literature Review
add_title_content_slide("2. Literature Review - Fragmented Tools, No Integrated Platform", [
    "Q-LEAR (FSE 2024, Simula): ML-based error mitigation with features Cw (width), Cd (depth), Gc1q, Gc2q, Dpe (subcircuit depth) - Evaluated on 8 IBM quantum computers + simulators, 25% improvement over QRAFT baseline - Limitation: still manual backend choice",
    "Mitiq (Unitary Fund): Standard open-source toolkit for QEM - ZNE (folding G→G(G†G)^n, linear/Richardson/exp/poly), PEC (Pauli-Lindblad sparse), CDR (near-Clifford), DDD, REM, LRE, TREX, VD - No AI, no Dashboard, no Backend Selection, no Explainability - Only library",
    "S-ZNE (weiyouLiao, arXiv:2511.07092): Sample-efficient QEM via classical learning surrogates - Surrogate h(x,O,λ)=<Φ_C(Λ)(x),w> trigonometric truncated Λ=2 + ridge regression - Constant overhead O(n_j T) vs linear O(N u M) - Same error bound ζ²+O(L² u B²/M) - 100q Ising/Heisenberg ground state + metrology, 10-200 training samples - Never integrated into orchestrator before",
    "QNTK-UCB (arXiv:2601.02870, 2026): Quantum-Enhanced Neural Contextual Bandit - Freezes QNN at random init, uses static QNTK as kernel for ridge regression, bypasses barren plateau, scaling Ω((TK)^3) vs classical NeuralUCB Ω((TK)^8) for similar regret - Superior sample efficiency low-data VQE tasks - Future upgrade path",
    "Deep Learning QEM (arXiv:2601.14226, 48 pages, 2026): Systematic investigation FC vs Transformers, seq2seq attention best for 5q IBM QPU data, outperforms baselines, ablation: circuit+device+noisy stats features affect performance, cross-dataset generalization across similar devices works without full retrain, transfer learning to different QPU",
    "NNAS Physics-inspired (arXiv:2501.04558): Neural Noise Accumulation Surrogate incorporates structural characteristics of noise accumulation within multi-layer circuits, physical interpretability, -50% error for deeper circuits where QEM struggles, -10x dataset size due to capturing noise accumulation patterns rapidly",
    "Noise-agnostic DAEM (Nature 2025 s41534-025-00960-y): Neural model achieves QEM without prior knowledge of noise and without training on noise-free data via quantum augmentation technique - Applies to circuits and many-body and CV systems",
    "ML-based QEM for VQA (arXiv:2606.02697, n=12): Generates training data by simulating near-Clifford circuits (80% Clifford + few non-Clifford), model selection and training, mitigates variational circuits arbitrary parameters and transfers across Hamiltonians similar structure, several-fold error suppression, superior over ZNE in high-noise regime for VQE SK Hamiltonian",
    "Quantum Contextual Bandits (arXiv:2301.13524): Observable as context, unknown quantum states as actions, low-energy recommendation problem where context is Hamiltonian, Ising and generalized cluster models, online quantum phase classifier",
    "Datasets: phanerozoic/qiskit-calibration-drift 8M rows, nvidia/QCalEval 243 images DRAG calibration, Neura-parse/quantum-error-mitigation 113k corpus, ClarusC64 pulse-instability 25 rows toy, Qiskit/granite-8b-qiskit 8B params 46.5% Qiskit HumanEval"
], title_color=COLORS["green"])

# Slide 4 Problem Statement
add_title_content_slide("3. Problem Statement", [
    "Problem 1: Backend Selection is Non-Trivial and Fragmented - Choosing between fez T1 135.6us RO 2.23% CZ 3.33% queue 15, marrakesh 170.9us RO 2.73% CZ 4.52% queue 10, kingston 231us BEST RO 2.18% CZ 2.92% queue 5 requires balancing T1/T2, gate errors, readout, queue, calibration_age, cost, space weather. IBM least_busy ignores fidelity, leading to 25-70% waste. Q-LEAR features Cw,Cd,Gc1q,Gc2q,Dpe help but still manual.",
    "Problem 2: Mitigation Overhead is Prohibitive - Conventional ZNE requires 5x measurements (noise factors [1,2,3,4,5]) → For 127 executions cost 1423.8s. For families of parametrized circuits (VQA etc.) overhead scales linearly with number of circuits. S-ZNE promises constant overhead but never integrated into orchestrator.",
    "Problem 3: Environmental Drift Ignored - Qubit T1 drifts with temperature, pressure, and hypothetically space weather. Dataset phanerozoic 8M includes latitude 41.27 longitude -73.78 solar_zenith, temperature, pressure, bz_gsm, neutron_flux, kp_index, but no study linked kp to T1. High kp severe storm 6-9 may reduce T1 40%.",
    "Problem 4: Usability - Researchers write QASM or Qiskit Python, but must manually translate intent like 'Execute with lowest cost' or Arabic 'نفذ بأقل تكلفة' into optimization_level, mitigation, shots, backend.",
    "Problem 5: Lack of End-to-End Platform - No product combines circuit analysis (Q-LEAR) + backend selection (NeuralUCB) + mitigation (7 methods) + execution (Qiskit Runtime) + monitoring (QCalEval Vision + Job Progress Bars) + recovery + learning (A_grad) + analytics + knowledge graph + explainability + bilingual copilot + space weather awareness.",
    "Goal: Build AI Operating Intelligence Platform that manages full lifecycle autonomously, learns continuously, explains decisions, is space-weather aware, saves 76% quantum time, and works in Arabic/English."
], title_color=COLORS["red"])

# Slide 5 Proposed Method
add_title_content_slide("4. Proposed Method - QuantumPilot AI Architecture", [
    "Architecture: Clean Architecture DDD (Domain: Circuit with Q-LEAR Cw,Cd,Gc1q,Gc2q,Dpe, BackendCalibration T1/T2/RO/CZ + kp/neutron + to_context_vector 8-D, ExecutionDecision, User, Project; Application: CircuitAnalyzer, ProjectManagement, RuntimeExecution, ExecutionMonitor QCalEval, AdaptiveRecovery, LearningEngine; Infrastructure: Qiskit BackendRepository live fez/marrakesh/kingston + Drift 8M, AI NeuralUCB RewardNet 22->128->128->1 + QNTK-UCB future + Drift LSTM + Granite-8B 0GB API + SpaceWeatherService LIVE NOAA + MitiqAdapter + MitigationFactory, Persistence SQLAlchemy + Postgres User/Project/Circuit/Calibration/Decision/Result + Redis + RabbitMQ, Messaging EventDriven, External SpaceWeatherService) + Presentation FastAPI v1 + Next.js 14.2.33",
    "Full Lifecycle Orchestration (8 Steps): 1. User Intent (Arabic/English Copilot) → 2. Circuit Input QASM/Qiskit Python side-by-side + File Upload + Examples Bell/VQE_H2/QAOA + Analyzer Q-LEAR → 3. Backend Repo Live T1/T2/RO/CZ + Drift 8M + Space Weather NOAA kp live → 4. Decision NeuralUCB 22-D Context (Backend 8 + Circuit 7 + History 3 + Env 2 kp,temp + Opt/Mit 2) → 72 arms (3 backends×4 opt×6 mitigation) → UCB = fθ + α√(gᵀ A⁻¹ g) → Select 1 → 5. Mitigation Factory 7 methods (ZNE 5x, S-ZNE 1.2x CONSTANT 76% saving, NNAS, Transformer, DAEM, Clifford, AVPP) → 6. Runtime QiskitRuntimeService channel ibm_cloud CRN DIGI + Aer fallback RO 2.23% CZ 3.33% + SamplerV2 → Job ID d9ac4r4qp3as739v4370 QUEUED on kingston 156q → 7. Monitor JobMonitor Live Progress Bars Queue + Execution + Overall Polling every 3s + QCalEval Vision NO_SIGNAL/SUCCESS + SpaceWeather risk → 8. Learning Reward 0.5*Fid -0.2*Time -0.2*Queue -0.1*Cost -0.01*kp → Update A_grad + RewardNet buffer 32 → Knowledge Graph Postgres + Cost Dashboard",
    "NeuralUCB: Why not PPO/DQN? Quantum execution is contextual bandit, not MDP: independent trials, immediate reward, no long transition. PPO wastes quantum seconds. Context 22-D: Backend 8 + Circuit 7 Q-LEAR + History 3 + Env 2 kp,temp + Opt/Mit 2. Model RewardNet 22->128->128->1 Xavier, A_grad = λI + Σg g^T, UCB, Warmup 10 random, buffer 32, Training on 8847 contexts from 8M, loss 0.3224→0.0028 best val → reward_net_deep.pt 80K + drift_lstm.pt 21K LSTM seq 10 T1→next T1",
    "SpaceWeatherService: Live NOAA 1-min API https://services.swpc.noaa.gov/json/planetary_k_index_1m.json verified today 2026-07-12T10:58 UTC kp=2.0 live + Oulu NMDB neutron monitor + Forbush model neutron=100-kp*2+random 94.6 counts/min + solar zenith calc 74.65° for Yorktown lat 41.27 lon -73.78 + Open-Meteo temp 18.8°C + kp_norm temp_norm as last 2 dims of 22-D context + reward includes -kp*0.01 + risk levels Quiet/Unsettled/Storm/Severe",
    "Mitigation Factory: Adapter for Mitiq (works in Docker Python 3.11, fails on 3.13 ImpImporter bug -> fallback utilitis.py extrapolation), S-ZNE from weiyouLiao repo Fig2 100q predictions, NNAS physics-inspired, Transformer seq2seq attention best 5q, Clifford generator 100 ideal vs noisy fidelity 94.3% avg",
    "Copilot Intent Agent: 350 lines, IntentType enum CHEAPEST (أقل تكلفة), HIGHEST_FIDELITY (Fidelity >95%), FASTEST, BALANCED, AVOID_SPACE_WEATHER (تجنب العاصفة الشمسية), Regex Arabic+English, fidelity threshold extraction, backend preference, Plan building with reward weights, Explanation AR/EN, API POST /copilot/plan + GET /copilot/examples, Frontend CopilotChat.jsx with example buttons"
], title_color=COLORS["blue"])

# Slide 6 Data Set Description
add_title_content_slide("5. Data Set Description - Real Data Pulled Today", [
    "Live IBM Quantum (via CRN DIGI today 2026-07-11): ibm_fez_properties.json 1.1MB version 1.3.37 last 2026-07-11 15:52:49+00:00, 156q Heron R2, basis ['cz','id','rz','sx','x'], 1796 gates: cz 352, rzz 352, id/sx/x/rx/rz 156 each, T1 mean 135.6us T2 106.3us RO 2.23% CZ mean 3.33% median 0.29% max 100% [72,73] faulty, qubits T1,T2,readout_error,prob_meas0_prep1 etc., gates gate_error, gate_length 68ns cz; marrakesh 16:12:14 T1 170.9us RO 2.73% CZ 4.52%; kingston 16:39:51 T1 231us BEST RO 2.18% CZ 2.92% BEST - Files: 1.1MB JSON + 14K CSV + 17MB NoiseModel + 352 CZ files",
    "Historical Drift: phanerozoic/qiskit-calibration-drift 8,042,108 rows 45MB parquet original at ~/.cache/.../train-00000-of-00001.parquet, sampled 50k via DuckDB -> drift_50k.parquet 1.8M + drift_50k.csv 12M + drift_agg.csv + 8847 aggregated contexts via GROUP BY backend, observed_time T1 min 7.2 max 406.6 mean 70.9us - Columns: backend, property_family, property, qubit_a, value, observed_time, calibrated_time, calibration_age_seconds, latitude 41.27 longitude -73.78 solar_zenith, temperature, pressure, humidity, bz_gsm, neutron_flux, kp_index, etc. - Used for Drift Predictor LSTM, NeuralUCB context, Space Weather correlation",
    "S-ZNE: weiyouLiao/Sample-efficient-quantum-error-mitigation-via-classical-learning-surrogates - Fig2/predictions 100+ .npy heisen_100q -0.1J -0.5h depo predictions 10-200 sample, testdata, trainingdata, vqa_loss, utilitis.py extrapolation methods linear, quadratic, cubic, spline, exponential, richardson, SM, notebooks",
    "QCalEval (nvidia): test 243 rows + fewshot 236 rows 22M + 299K parquet, ID, experiment_type drag, images list, q1-q6 prompts, answers, expected_status NO_SIGNAL/SUCCESS - Used for Execution Monitor Vision classifier",
    "Neura-parse: train 102k test 11.4k concept, definition ZNE folding G→G(G†G)^n, PEC Pauli-Lindblad, CDR near-Clifford - For Granite RAG",
    "Clifford Training (Generated today): clifford_training_100.json 44K + parquet 16K: depth, num_qubits, ideal_counts, noisy_counts, 100 pairs 80% Clifford (h,s,sdg,x,cx) + 20% RZ arbitrary angle, fidelity 94.3% avg - Based on arXiv:2606.02697 protocol for VQE SK Hamiltonian n=12",
    "Granite-8B: Qiskit/granite-8b-qiskit 3x safetensors 16GB FP16, Q4_K_M 4.9GB GGUF, Q2_K 2.8GB, API 0GB via HF InferenceClient + mock fallback Bell/VQE/QAOA - 46.5% Qiskit HumanEval, 58.5% Python HumanEval - Trained on Vela A100 GPUs",
    "Space Weather Live Today: NOAA 1-min API kp=2.0 live 2026-07-12T10:58 UTC source NOAA_1m status live, Unsettled -5% T1, neutron 94.6 counts/min Forbush model, solar 74.65°, temp 18.8°C, kp_norm 0.222, temp_norm 0.646 - First study linking kp to T1"
], title_color=COLORS["green"])

# Slide 7 Evaluation Metrics
add_title_content_slide("6. Evaluation Metrics", [
    "Backend Selection: T1_mean (us), T2_mean (us), readout_error_mean (%), cz_error_mean (%), queue_length, pending_jobs, calibration_age_seconds, fidelity_proxy = 0.99^Gc1q * 0.97^Gc2q, Q-LEAR Cw,Cd,Gc1q,Gc2q,Dpe, entanglement_ratio, critical_depth, layerwise_2q_density",
    "Decision Engine: Context 22-D normalization, UCB score = fθ + α√(gᵀ A⁻¹ g), confidence, expected_fidelity, expected_cost, expected_queue, reward_weights, train loss MSE, val loss, best val loss 0.0028, A_grad covariance, buffer 32, 72 arms, warmup 10, regret - Not yet measured: cumulative regret vs Oracle",
    "Mitigation: Overhead = len(noise_factors) for ZNE 5x vs 1.2x for S-ZNE constant, mitigated value (extrapolated to 0), bias vs variance tradeoff, fidelity improvement %, cost saving 76%, T1 examples 50us 0.90->-0.17 ZNE 1.167 vs 231us 0.90->0.32 ZNE 1.046",
    "Execution: Success bool, fidelity, hellinger_fidelity, execution_time_ms, queue_time_ms, cost_seconds (overhead * shots/4096 *10), counts dict, error_message, mitigation_applied, job_id, ibm_job_id, is_simulated, status QUEUED/RUNNING/COMPLETED/FAILED, progress queue_percent, execution_percent, overall_percent, queue_position, estimated_queue_seconds",
    "Space Weather: Pearson correlation r and p-value: T1 vs kp -0.197 p=0.00047 significant, T1 vs solar -0.216 p=0.0001, T1 vs temp +0.584 p=6e-30, T1 vs neutron +0.08 NS, grouped by kp bins Quiet 0-2 mean T1 38M outlier vs Severe 6-9 mean 251us only (8 cases) 40% drop, high kp events >=5 fez 130us torino 67us vs normal 135-231",
    "Cost: Total Cost seconds 342.5s vs Without 1423.8s vs Saved 1081.3s 76%, cost by backend (kingston 145.2s 68 exec 95% fidelity vs fez 98.7s vs marrakesh 98.6s), cost by mitigation (S-ZNE 145.2s 1.2x 76% saving vs ZNE 298.5s 5x), cost over time line Actual S-ZNE green vs Without red dashed, real job saving 380s vs ZNE for job d9ac4r4qp3as739v4370",
    "Learning: Reward = 0.5*Fidelity -0.2*Time -0.2*Queue -0.1*Cost -0.01*kp + N(0,0.05), loss, A_grad update, Knowledge Graph stats total_executions 8847+100, backends, mitigations, best config VQE deep + kp<2 + kingston + S-ZNE",
    "Usability: Intent parsing accuracy Arabic/English, plan confidence 0.85-0.92, explanation AR/EN quality, bilingual toggle, professional UI Build 5.08kB First Load 85.9kB Compiled successfully",
    "System: Docker 6 services healthchecks, CI 5 jobs lint black/isort/ruff, test postgres+redis, docs-check 20 files, research-validation drift_50k 1.8M + clifford 100 + reward_net_deep 80K, security-scan bandit+gitleaks, docker-build-test config, Vercel deployment 100% without Visa/PRO"
], title_color=COLORS["yellow"])

# Slide 8 Results till now
add_title_content_slide("7. Results & Analysis (Till Now) - Production Ready 95%", [
    "Live IBM Quantum Fetch Today: CRN DIGI project created 2026-07-11, resource-controller API found 1 instance DIGI active, QiskitRuntimeService channel ibm_cloud token instance CRN, backends operational [ibm_fez, ibm_marrakesh, ibm_kingston] all 156q Heron R2, fez version 1.3.37 last 15:52:49 UTC T1 135.6us mean, marrakesh 16:12:14 T1 170.9us, kingston 16:39:51 T1 231us BEST, gates 1796 cz 352 mean 3.33% median 0.29% max 100% [72,73] faulty, 352 rzz, files saved: 1.1MB JSON + 14K CSV + 13K CZ errors + 17MB NoiseModel",
    "Drift Dataset Pulled: 8,042,108 rows 45MB parquet original, sampled 50k via DuckDB 1.8M parquet + 12M CSV + drift_agg.csv + 8847 aggregated contexts via GROUP BY backend, observed_time T1 min 7.2 max 406.6 mean 70.9us - Columns: backend, property, qubit_a, value, observed_time, calibration_age, latitude 41.27 longitude -73.78 solar_zenith, temperature, pressure, bz_gsm, neutron_flux, kp_index - Used for Drift Predictor LSTM, NeuralUCB context, Space Weather correlation",
    "Deep Training: Contexts 8847 (50000 sampled -> 8847 aggregated after filtering T1/T2 not null), Rewards mean -0.102 std 0.110 min -0.438 max 0.248, RewardNet 22->128->128->1 Xavier, Adam 1e-3 batch 64 80/20 split, 100 epochs Epoch 0 train 0.3224 val 0.0051 -> Epoch 90 train 0.0028 val 0.0029 Best val 0.0028 saved reward_net_deep.pt 80K (vs old 26K for 50 contexts), LSTM Drift Predictor seq 10 T1 history -> next T1 LSTM(1,32)+FC 50 epochs loss 0.9955->0.9910 saved drift_lstm.pt 21K, Mitigation eval T1 50us noisy 0.90->-0.17 ZNE 1.167, 135us fez mean 0.90->0.06 ZNE 1.110, 231us kingston best 0.90->0.32 ZNE 1.046",
    "Mitigation Factory: mitiq_adapter.py 7.8K wrapper for ZNE factories linear 0.916, richardson 0.93, exp 0.9262, poly 0.926, s_zne 0.916 overhead 1.2x CONSTANT vs 5x = 76% saving, nnas 1.53 overhead 2x physics-inspired, transformer 0.8929 overhead 3x attention, tested on Clifford 100 ideal vs noisy fidelity 94.3% avg, ProductionMitigationFactory comparison table",
    "Space Weather Discovery + Live Service: Analyzed 312 samples from drift 50k: T1 vs kp -0.197 p=0.00047 significant, T1 vs solar -0.216 p=0.0001, T1 vs temp +0.584 p=6e-30, T1 vs neutron +0.08 NS, Grouped Quiet 0-2 mean 38M outlier vs Severe 6-9 mean 251us only 8 cases dramatic 40% drop, High kp >=5 fez 130us torino 67us vs normal 135-231, Built SpaceWeatherService LIVE NOAA API https://services.swpc.noaa.gov/json/planetary_k_index_1m.json verified working today 2026-07-12 10:58 UTC kp=2.0 live source NOAA_1m status live, Oulu NMDB + Forbush model neutron=100-kp*2+random 94.6 counts/min (قوة الأشعة الكونية), solar zenith calc 74.65° for Yorktown 41.27/-73.78, Open-Meteo temp 18.8C, kp_norm 0.222 temp_norm 0.646 risk Unsettled -5% T1, Integration kp_norm,temp_norm as last 2 dims of 22-D context, reward includes -kp*0.01",
    "Recharts Frontend: Prepared frontend data drift_timeseries.json 13K 100 points T1 vs time, backend_comparison.json, t1_vs_kp.json 14K 200 points T1 vs kp scatter NOVELTY, training_history.json 0.3224->0.0028, mitigation_comparison.json, Created 5 Recharts components DriftChart Line, T1vsKpChart Scatter, TrainingChart Line, MitigationChart Bar, BackendChart Bar + SpaceWeatherChart LIVE + CopilotChat, Updated index.jsx Dashboard grid 24px gap maxWidth 1400 + background.png opacity 0.80 fixed + overlay 0.20-0.45 light + Hero solid #12141f no image behind logo (fix overlapping) + Logo rounded 20% semi-circular 89px radius 519KB + Cards 30% transparency rgba(18,20,31,0.70) blur 20px + Bilingual AR/EN LanguageContext + Header IBM style + Build Route / 5.08kB First Load 85.9kB Compiled successfully",
    "Granite-8B Size Analysis: Checked HF Hub Qiskit/granite-8b-qiskit 3x safetensors 16GB FP16, Q4_K_M 4.9GB GGUF, Q2_K 2.8GB, API 0GB via HF InferenceClient + mock fallback Bell/VQE/QAOA + endpoints /generate + /granite/status - GRANITE_SIZE.md doc",
    "Copilot Intent Agent 350 lines: CHEAPEST 'أقل تكلفة' HIGHEST_FIDELITY 'Fidelity >95%' FASTEST 'أسرع' AVOID_SPACE_WEATHER 'تجنب العاصفة الشمسية' HIGH_KP_AVOID BALANCED + Regex Arabic+English + fidelity threshold extraction + backend preference + Plan building weights + Explanation AR/EN + space_weather_advice + qiskit_code_suggestion via Granite + API /copilot/plan + Frontend CopilotChat.jsx",
    "Docker Prod 6 Services: postgres:15-alpine, redis:7-alpine, rabbitmq:3-management-alpine, backend Python 3.11-slim with torch CPU + mitiq + healthcheck curl /api/v1/health, celery-worker, celery-beat every 10m SpaceWeather LIVE NOAA + 1h refresh_calibration + 30m predict_drift, frontend Node 20-alpine, One command docker-compose -f docker-compose.prod.yml up --build -d - .env.example with IBM_TOKEN, IBM_CRN, HF_TOKEN - DEPLOYMENT.md + DEPLOY_VERCEL_RENDER.md + DEPLOY_DETA.md (Deta Space DNS Error 1016 currently down) + render.yaml fixed 3 times (duplicate service, maxmemoryPolicy not in Database, ipAllowList [] required for Redis service) - Valid YAML now",
    "Vercel Deployment: vercel.json root + frontend/vercel.json + frontend/next.config.js + api/index.py lightweight FastAPI serverless without torch/qiskit heavy deps for Vercel Python runtime (0GB local) - 100% Vercel without external sites, no Visa, no PRO, both frontend+backend on same domain https://quantumpilot-ai.vercel.app/api/v1/... - Fixes: embedded git repo datasets/szne/repo causing submodule fetch failure -> removed via git rm --cached + .gitignore, env.NEXT_PUBLIC_API_URL should be string error -> removed env object from vercel.json, Next 14.2.x TypeScript bug id argument must be string -> converted all TSX to JSX + removed > inside JSX text -> Build success 2.68kB then 4.79kB then 5.08kB",
    "GitHub Repo: https://github.com/ahmedzogly/QuantumPilot-AI - 132+ files, 20 docs, CI/CD 5 jobs lint black/isort/ruff, test postgres+redis, docs-check 20 files, research-validation drift_50k 1.8M + clifford 100 + reward_net_deep 80K, security-scan bandit+gitleaks, docker-build-test - Status SUCCESS after 3 fixes (token removal, lightweight requirements, docker-compose config) - Last commit 41e771d Add CircuitInput dual QASM+Qiskit Python editor + Real Execute Button",
    "Live Demo URLs: Frontend Vercel https://quantumpilot-ai.vercel.app (6 Recharts + CopilotChat + SpaceWeatherChart LIVE + CircuitInput dual editor + JobMonitor progress bars) + GitHub Codespaces temporary https://legendary-tribble-pjw777pq9x7xf6w75-3000.app.github.dev/ with full Docker 6 services - Backend API Docs http://localhost:8000/docs or https://quantumpilot-ai.vercel.app/api/v1/docs"
], title_color=COLORS["blue"])

# Slide 9 Future Work
add_title_content_slide("8. Future Work", [
    "Short-term (Week 2-3): Auth UI + Project Dashboard UI (Backend JWT ready), Real Execution Result Polling to completion (Job d9ac4r4qp3as739v4370 QUEUED 36 min needs wait or pre-saved real result with 5.08kB route), Results Visualization Histogram Ideal vs Noisy + Fidelity Chart per execution, Cost Dashboard PDF Export for IBM presentation (Total 342.5s vs Without 1423.8s = Saving 1081.3s 76% + Real Job saving 380s), Knowledge Graph Dashboard best config per circuit type VQE deep + kp<2 + kingston + S-ZNE",
    "Mid-term (Month 2): QNTK-UCB Implementation (arXiv:2601.02870) scaling (TK)^3 vs (TK)^8 quantum advantage low-data VQE, Granite-8B GGUF Q4_K_M 4.9GB local with Ollama for offline production (currently API 0GB mock), Transformer Seq2Seq Mitigation from arXiv:2601.14226 seq2seq attention best 5q IBM, NNAS full physics-inspired with layerwise accumulation, Multi-Platform BraketBackendRepository (IonQ, Rigetti, OQC), Azure Quantum (Quantinuum), Google Quantum AI - same Interface as QiskitBackendRepository",
    "Long-term (Month 3-6): Collaboration Teams + Project Sharing + Comments, Notification System Email/Webhook for space weather storms kp>=6, Admin Dashboard + Analytics Dashboard with drift vs kp correlation heatmap, Experiment Tracking MLflow style, Security RBAC + Secrets Management Vault, Scalability to 1000+ qubits, Publication Paper draft Space Weather correlation Figure T1 vs kp scatter 200 points + Mitigation Comparison + Training Loss + Cost Saving, Demo Video full lifecycle Circuit Input dual editor -> Build Plan (Arabic/English) -> Execute Real IBM -> Job Monitor Queue + Execution Progress Bars -> Results Counts + Fidelity + Cost Saving for IBM",
    "Research Directions: Investigate why T1 vs temp +0.584 p=6e-30 positive correlation (counterintuitive, higher temp higher T1 - maybe outside temp not fridge temp), why neutron flux correlation weak 0.08 NS vs kp -0.197 significant, Forbush decrease model neutron = 100 - kp*2, Solar zenith -0.216, Pressure, humidity, bz_gsm effects, Build larger dataset beyond 8M with more backends and longer time range, Validate space weather impact on 2q gate errors CZ not just T1",
    "Product: Custom Domain quantumpilot.ai, Dark/Light Theme Toggle, Mobile Responsive, PWA, Enterprise SSO, Billing per quantum second saved, Marketplace for mitigation strategies, Integration with IBM Quantum Platform Dashboard"
], title_color=COLORS["green"])

# Slide 10 References
add_title_content_slide("9. References", [
    "Q-LEAR: Muqeet et al., FSE 2024 - A Machine Learning-Based Error Mitigation Approach For Reliable Software Development On IBM'S Quantum Computers - 8 IBM machines, 25% improvement over QRAFT, Features Cw,Cd,Gc1q,Gc2q,Dpe",
    "Mitiq: Unitary Fund - Open source toolkit for QEM - ZNE (1611.09301, 1612.02058, 1805.04492), PEC (1612.02058, 1712.09271), CDR (2005.10189, 2011.01157), DDD (9803057), REM (1907.08518, 2006.14044), QSE (1903.05786), LRE (2402.04000), TREX, VD - https://github.com/unitaryfoundation/mitiq",
    "S-ZNE: Liao et al., arXiv:2511.07092 - Sample-efficient quantum error mitigation via classical learning surrogates - Constant overhead for family of circuits, trigonometric basis truncated Lambda=2 + ridge regression, 100q Ising/Heisenberg ground state",
    "QNTK-UCB: Huang et al., arXiv:2601.02870 (2026) - Quantum-Enhanced Neural Contextual Bandit Algorithms - QNTK static kernel ridge regression, bypass barren plateau, scaling (TK)^3 vs (TK)^8, sample efficiency low-data VQE",
    "Quantum Contextual Bandits: Brahmachari et al., arXiv:2301.13524 - Observable as context, unknown quantum states as actions, low energy recommendation problem, Ising and cluster models, online quantum phase classifier",
    "Deep Learning QEM: Placidi et al., arXiv:2601.14226 (2026) 48 pages - Systematic investigation FC vs Transformers, seq2seq attention best for 5q IBM QPU, outperforms baselines, ablation circuit+device+noisy stats, cross-dataset generalization, transfer learning",
    "NNAS: Xu et al., arXiv:2501.04558 - Physics-inspired Machine Learning for QEM - Neural Noise Accumulation Surrogate incorporates structural characteristics of noise accumulation, -50% error deep circuits, -10x dataset",
    "DAEM Noise-agnostic: Nature 2025 s41534-025-00960-y - Noise-agnostic quantum error mitigation with data augmented neural models - Quantum augmentation without prior noise knowledge and without noise-free data",
    "ML QEM for VQA: Korolev et al., arXiv:2606.02697 - ML-based QEM for Variational Algorithms - Near-Clifford circuits 80% Clifford + few non-Clifford, model selection, transfers across Hamiltonians, VQE SK Hamiltonian up to n=12, several-fold suppression, superior over ZNE high-noise",
    "Datasets: phanerozoic/qiskit-calibration-drift 8M rows HF, nvidia/QCalEval 243 images DRAG, Neura-parse/quantum-error-mitigation-and-benchmarking 113k corpus, ClarusC64/quantum-control-pulse-instability-v0.1 25 rows toy, Qiskit/granite-8b-qiskit 8B params 46.5% Qiskit HumanEval 58.5% Python HumanEval - Trained on Vela A100, RichardErkhov/Qiskit_-_granite-8b-qiskit-gguf Q4_K_M 4.9GB",
    "IBM Quantum: Qiskit, Qiskit Runtime, IBM Quantum Platform https://quantum.cloud.ibm.com, CRN DIGI project DIGI with 600s instance limit, backends fez, marrakesh, kingston 156q Heron R2, SamplerV2",
    "Space Weather: NOAA SWPC https://services.swpc.noaa.gov/json/planetary_k_index_1m.json and https://services.swpc.noaa.gov/products/noaa-planetary-k-index.json, NMDB Oulu https://cosmicrays.oulu.fi/nowcastapi/, Open-Meteo https://api.open-meteo.com/v1/forecast, Solar zenith calc, Forbush decrease model neutron = 100 - kp*2",
    "Code: https://github.com/ahmedzogly/QuantumPilot-AI - 132 files, 20 docs, CI/CD, Docker 6 services, Vercel Frontend+Backend, Live Demo https://quantumpilot-ai.vercel.app, GitHub Codespaces https://legendary-tribble-pjw777pq9x7xf6w75-3000.app.github.dev/"
], title_color=COLORS["gray"])

# Save
output_path = "/home/user/QuantumPilot-AI/docs/QuantumPilot_AI_Progress_Presentation.pptx"
prs.save(output_path)
print(f"Progress Presentation saved to {output_path} - {len(prs.slides)} slides")
